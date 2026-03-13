from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt as PPTPt
import tempfile
import os
import shutil

app = FastAPI()


@app.post("/download-doc")
async def download_doc(
    file: UploadFile = File(...),
    document_code: str = Form(""),
    client_name: str = Form(""),
    department: str = Form(""),
    document_type: str = Form(""),
    purpose: str = Form(""),
    created_on: str = Form(""),
    created_by: str = Form("")
):

    try:

        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, file.filename)

        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        ext = os.path.splitext(file.filename)[1].lower()

        # ===============================
        # WORD DOCUMENT PROCESSING
        # ===============================
        if ext == ".docx":

            doc = Document(input_path)

            body = doc._element.body

            # PAGE BREAK
            page_break = OxmlElement("w:p")
            run = OxmlElement("w:r")
            br = OxmlElement("w:br")
            br.set(qn("w:type"), "page")
            run.append(br)
            page_break.append(run)

            # HEADING
            heading_para = doc.add_paragraph()
            heading_para.text = "DOCUMENT DETAILS"

            run = heading_para.runs[0]
            run.bold = True
            run.underline = True
            run.font.name = "Arial"
            run.font.size = Pt(22)

            heading_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # DETAILS
            details_para = doc.add_paragraph()
            details_para.paragraph_format.line_spacing = 2

            details_list = [
                f"Document Code : {document_code}",
                f"Client Name : {client_name}",
                f"Department : {department}",
                f"Document Type : {document_type}",
                f"Purpose : {purpose}",
                f"Created On : {created_on}",
                f"Created By : {created_by}",
            ]

            for d in details_list:
                r = details_para.add_run(d + "\n")
                r.font.name = "Arial"
                r.font.size = Pt(16)

            # INSERT AT TOP
            body.insert(0, page_break)
            body.insert(0, details_para._p)
            body.insert(0, heading_para._p)

            output_path = os.path.join(temp_dir, file.filename)
            doc.save(output_path)

        # ===============================
        # POWERPOINT PROCESSING
        # ===============================
        elif ext == ".pptx":

            prs = Presentation(input_path)

            # --------------------------------
            # FIND SAFE LAYOUT (TITLE + BODY)
            # --------------------------------
            chosen_layout = None

            for layout in prs.slide_layouts:

                title_found = False
                body_found = False

                for placeholder in layout.placeholders:

                    ph_type = placeholder.placeholder_format.type

                    if ph_type == 1:
                        title_found = True

                    if ph_type == 2:
                        body_found = True

                if title_found and body_found:
                    chosen_layout = layout
                    break

            # fallback
            if not chosen_layout:
                chosen_layout = prs.slide_layouts[0]

            detail_slide = prs.slides.add_slide(chosen_layout)

            # --------------------------------
            # MOVE SLIDE TO POSITION 2
            # --------------------------------
            slide_ids = prs.slides._sldIdLst
            slides = list(slide_ids)

            slide_ids.remove(slides[-1])
            slide_ids.insert(1, slides[-1])

            # --------------------------------
            # FIND TITLE
            # --------------------------------
            title_shape = None

            for shape in detail_slide.shapes:

                if shape.has_text_frame and shape == detail_slide.shapes.title:
                    title_shape = shape
                    break

            if title_shape:

                title_shape.text = "DOCUMENT DETAILS"

                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.size = PPTPt(40)

            # --------------------------------
            # FIND BODY PLACEHOLDER
            # --------------------------------
            body_shape = None

            for shape in detail_slide.shapes:

                if shape.has_text_frame and shape != title_shape:
                    body_shape = shape
                    break

            if body_shape:

                tf = body_shape.text_frame
                tf.clear()

                details = [
                    f"Document Code : {document_code}",
                    f"Client Name : {client_name}",
                    f"Department : {department}",
                    f"Document Type : {document_type}",
                    f"Purpose : {purpose}",
                    f"Created On : {created_on}",
                    f"Created By : {created_by}",
                ]

                for i, d in enumerate(details):

                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = d
                    p.level = 0

                    for run in p.runs:
                        run.font.name = "Arial"
                        run.font.size = PPTPt(24)

            output_path = os.path.join(temp_dir, file.filename)
            prs.save(output_path)

        else:
            raise HTTPException(
                status_code=400,
                detail="Only DOCX and PPTX files are supported"
            )

        return FileResponse(
            output_path,
            media_type="application/octet-stream",
            headers={
                "Content-Disposition": f'attachment; filename="{file.filename}"'
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
